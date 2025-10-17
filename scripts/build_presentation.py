# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np

# Create output and assets directories
os.makedirs("output", exist_ok=True)
os.makedirs("assets", exist_ok=True)

# Color palette
BLUE = RGBColor(30, 76, 120)
GREEN = RGBColor(88, 142, 38)
OLIVE = RGBColor(120, 158, 73)

# Slide size widescreen
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

# ---------- Helpers ----------
def add_title_shape(slide, title):
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12.0), Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.RIGHT  # RTL-friendly
    return tx

def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(slide, title)
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.4))
    tf = tx.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        para = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        para.text = b
        para.font.size = Pt(24)
        para.alignment = PP_ALIGN.RIGHT
    return slide

def add_image(slide, path, left=0.7, top=1.6, width=11.5):
    slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))

def draw_banner(slide, color=OLIVE):
    # Optional top banner for aesthetic
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.shadow.inherit = False

# ---------- Charts & Art generators ----------
def save_fig(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=220)
    plt.close(fig)

def build_prevalence_chart():
    # Data inspired by IDF Diabetes Atlas (global adults with diabetes, M)
    years = [2000, 2005, 2010, 2015, 2021, 2025]
    world = [151, 194, 226, 342, 537, 643]
    fig, ax = plt.subplots(figsize=(8,3.5))
    ax.plot(years, world, marker="o", color="#1f77b4", linewidth=2)
    ax.set_title("روند جهانی دیابت (تقریبی، منبع: IDF Diabetes Atlas)")
    ax.set_xlabel("سال")
    ax.set_ylabel("میلیون نفر")
    ax.grid(alpha=0.3)
    save_fig(fig, "assets/prevalence.png")

def build_hba1c_chart():
    labels = ["Esposito 2009", "DIRECT 2008\n(زیرگروه T2D)", "Meta-analysis 2013"]
    delta = [-0.9, -0.5, -0.4]  # تغییر HbA1c نسبت به مقایسه (تقریبی)
    fig, ax = plt.subplots(figsize=(8,3.5))
    bars = ax.bar(labels, delta, color="#2ca02c")
    ax.axhline(0, color="k", linewidth=0.8)
    ax.set_ylabel("تغییر HbA1c (%) – مدیترانه‌ای نسبت به مقایسه")
    for r in bars:
        ax.text(r.get_x()+r.get_width()/2, r.get_height()+(-0.05), f"{r.get_height():.1f}",
                ha="center", va="top", color="white", fontsize=10)
    save_fig(fig, "assets/hba1c.png")

def build_inflammation_chart():
    markers = ["CRP", "IL-6", "TNF-α"]
    change = [-25, -12, -10]  # درصد کاهش تقریبی
    fig, ax = plt.subplots(figsize=(8,3.5))
    bars = ax.bar(markers, change, color="#8abf69")
    ax.axhline(0, color="k", linewidth=0.8)
    ax.set_ylabel("تغییر (%)")
    ax.set_title("کاهش مارکرهای التهابی با الگوی مدیترانه‌ای (تقریبی)")
    for r in bars:
        ax.text(r.get_x()+r.get_width()/2, r.get_height()+(-1.5), f"{int(r.get_height())}%",
                ha="center", va="top", color="white", fontsize=10)
    save_fig(fig, "assets/inflammation.png")

def build_med_pyramid():
    # Simple pyramid schematic
    fig, ax = plt.subplots(figsize=(6,4))
    ax.axis("off")
    levels = ["سبزیجات/میوه‌ها/غلات کامل/حبوبات", "روغن زیتون/مغزدانه‌ها", "ماهی/لبنیات", "گوشت قرمز/شیرینی‌ها (کم)"]
    colors = ["#e0f2f1", "#dcedc8", "#fff9c4", "#ffe0b2"]
    for i, (lvl, c) in enumerate(zip(levels, colors)):
        pass
    # Draw levels
    fig, ax = plt.subplots(figsize=(6,4))
    ax.axis("off")
    for i, (lvl, c) in enumerate(zip(levels, colors)):
        ax.add_patch(plt.Rectangle((0.1+i*0.05, 0.1+i*0.18), 0.8-0.1*i, 0.16, color=c, ec="#888"))
        ax.text(0.5, 0.18+i*0.18, lvl, ha="center", va="center", fontsize=10)
    ax.set_title("هرم ساده رژیم مدیترانه‌ای")
    save_fig(fig, "assets/med_pyramid.png")

def build_mito_art():
    fig, ax = plt.subplots(figsize=(6,3.5))
    ax.axis("off")
    # mitochondria oval
    oval = plt.Circle((0.5, 0.5), 0.45, color="#ffe0b2", ec="#f57c00")
    ax.add_patch(oval)
    for x in np.linspace(0.2, 0.8, 6):
        ax.plot([x, x], [0.2, 0.8], color="#f57c00", linewidth=2)
    ax.text(0.5, 0.1, "محافظت میتوکندری (ROS↓, PGC-1α↑)", ha="center", fontsize=11)
    save_fig(fig, "assets/mito.png")

def build_nfkb_art():
    fig, ax = plt.subplots(figsize=(6,3.5))
    ax.axis("off")
    ax.text(0.5, 0.8, "NF-κB / NLRP3", ha="center", fontsize=12, weight="bold")
    ax.arrow(0.2, 0.6, 0.6, 0, head_width=0.03, head_length=0.02, color="#c62828")
    ax.text(0.5, 0.62, "التهاب سیستمیک", ha="center", color="#c62828")
    ax.arrow(0.5, 0.55, 0, -0.25, head_width=0.03, head_length=0.02, color="#2e7d32")
    ax.text(0.52, 0.4, "مدیترانه‌ای → مهار مسیر", ha="left", color="#2e7d32")
    save_fig(fig, "assets/nfkb.png")

def build_gut_immune_art():
    fig, ax = plt.subplots(figsize=(6,3.5))
    ax.axis("off")
    ax.text(0.5, 0.8, "میکروبیوم روده → SCFA (بوتیرات)", ha="center", fontsize=12)
    ax.arrow(0.2, 0.6, 0.6, 0, head_width=0.03, head_length=0.02, color="#1565c0")
    ax.text(0.5, 0.62, "GLP-1 ↑  /  حساسیت به انسولین ↑", ha="center", color="#1565c0")
    ax.arrow(0.5, 0.55, 0, -0.25, head_width=0.03, head_length=0.02, color="#2e7d32")
    ax.text(0.52, 0.4, "التهاب مزمن ↓", ha="left", color="#2e7d32")
    save_fig(fig, "assets/gut_immune.png")

# ---------- Tables ----------
def add_sources_table(slide_title, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(slide, slide_title)
    cols = ["عنوان مقاله","نوع مطالعه/مدل","جمعیت/نمونه","هدف مطالعه","روش‌ها/طراحی","نتایج کلیدی","نتیجه‌گیری مرتبط با T2D","محدودیت‌ها/نکات مهم"]
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(n_rows, len(cols), Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.7)).table

    # Header
    for j, c in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = c
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.RIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN

    # Rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.RIGHT

    return slide

# ---------- Slides ----------
def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    draw_banner(slide, color=OLIVE)
    # Title
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(2.0))
    tf = tx.text_frame
    tf.text = "تأثیر رژیم غذایی مدیترانه‌ای بر دیابت نوع 2"
    p0 = tf.paragraphs[0]
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = BLUE
    p0.alignment = PP_ALIGN.RIGHT
    # Subtitle
    p1 = tf.add_paragraph()
    p1.text = "ارائه برای رشته علوم تغذیه"
    p1.font.size = Pt(24)
    p1.alignment = PP_ALIGN.RIGHT

    # Optional logo
    logo_path = "assets/logo.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.7), Inches(0.7), height=Inches(1.0))

def build_all_assets():
    build_prevalence_chart()
    build_hba1c_chart()
    build_inflammation_chart()
    build_med_pyramid()
    build_mito_art()
    build_nfkb_art()
    build_gut_immune_art()

def build_deck():
    build_all_assets()

    # 1
    add_title_slide()

    # 2
    add_bullets_slide("تعریف و ماهیت دیابت نوع 2", [
        "اختلال مزمن متابولیک با هیپرگلیسمی ناشی از مقاومت به انسولین و نقص ترشح آن",
        "درگیری کبد، عضله و بافت چربی؛ التهاب مزمن خفیف و استرس اکسیداتیو",
        "عوارض: قلبی‌–عروقی، نفروپاتی، نوروپاتی، رتینوپاتی، کبد چرب غیرالکلی"
    ])

    # 3
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "اپیدمیولوژی و شیوع")
    add_image(s, "assets/prevalence.png", left=0.9, top=1.8, width=10.8)

    # 4
    add_bullets_slide("عوامل خطر و تفاوت‌های جمعیتی", [
        "سن، سابقه خانوادگی، چاقی مرکزی، کم‌تحرکی",
        "الگوی غذایی پرکالری/فراوری‌شده، وضعیت اجتماعی–اقتصادی",
        "تفاوت‌های منطقه‌ای و جمعیتی در شیوع"
    ])

    # 5
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "روند شیوع در سال‌های اخیر")
    add_image(s, "assets/prevalence.png", left=0.9, top=1.8, width=10.8)

    # 6
    add_bullets_slide("درمان‌های موجود برای T2D", [
        "اصلاح سبک زندگی: تغذیه، فعالیت بدنی، خواب",
        "داروها: متفورمین، SGLT2i، GLP-1 RA، انسولین",
        "محدودیت‌ها: عوارض، هزینه، پایبندی — نقش کلیدی تغذیه"
    ])

    # 7
    add_bullets_slide("درمان‌های مکمل", [
        "مشاوره و آموزش تغذیه‌ای، رفتاردرمانی، مدیریت استرس",
        "بررسی مکمل‌ها با احتیاط و شواهد",
        "الگوهای اثربخش: مدیترانه‌ای، DASH، کم‌کربوهیدرات (انتخاب فردمحور)"
    ])

    # 8
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "رژیم مدیترانه‌ای: سوخت متفاوت برای بدن")
    add_image(s, "assets/med_pyramid.png", left=3.0, top=1.8, width=7.5)

    # 9
    add_bullets_slide("کاربردهای شناخته‌شده رژیم مدیترانه‌ای", [
        "کاهش ریسک CVD، بهبود فشارخون و لیپیدها",
        "بهبود NAFLD، برخی سرطان‌ها، سلامت شناخت و خلق",
        "کاهش خطر دیابت و بهبود کنترل قند در مبتلایان"
    ])

    # 10
    add_bullets_slide("چرا مدیترانه‌ای برای T2D مفید است؟", [
        "فیبر بالا → کاهش شاخص گلیسمی و افزایش سیری",
        "MUFA/PUFA → بهبود حساسیت به انسولین",
        "پلی‌فنول‌ها → ضدالتهاب و آنتی‌اکسیدان",
        "تنوع و انعطاف‌پذیری → پایبندی بهتر"
    ])

    # 11
    s = add_bullets_slide("مقدمه مکانیسم‌ها", [
        "مسیرها: التهاب، میتوکندری، گلوتامات/تحریک‌سمّیت، ایمنی، میکروبیوم، هورمون‌های روده",
        "از جزء غذایی → تغییرات مولکولی/سیستمی → پیامدهای بالینی"
    ])
    draw_banner(s)

    # 12
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "1) کاهش التهاب (Neuroinflammation ↓)")
    add_image(s, "assets/nfkb.png", left=3.2, top=1.8, width=7.0)

    # 13
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "نتیجه کاهش التهاب")
    add_image(s, "assets/inflammation.png", left=1.0, top=1.8, width=10.5)

    # 14
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "2) محافظت میتوکندریایی و بهبود انرژی")
    add_image(s, "assets/mito.png", left=3.0, top=1.8, width=7.5)

    # 15
    add_bullets_slide("آسیب میتوکندری در T2D", [
        "افزایش ROS و افت تولید ATP",
        "اختلال عملکرد سلول β و لیپوتوکسیسیته/گلوکوتوکسیسیته",
        "ارتباط با مقاومت به انسولین و پیشرفت عوارض"
    ])

    # 16
    add_bullets_slide("متابولیسم در رژیم مدیترانه‌ای", [
        "کنترل بار گلیسمی وعده‌ها و کاهش واریانس پس‌غذایی",
        "بهبود پروفایل لیپیدی و التهاب سیستمیک",
        "پایداری انرژی و سیری"
    ])

    # 17
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "اثرات محافظتی بر شاخص‌های بالینی")
    add_image(s, "assets/hba1c.png", left=1.0, top=1.8, width=10.5)

    # 18
    add_bullets_slide("3) تنظیم سیستم گلوتامات و کاهش Excitotoxicity", [
        "در عوارض عصبی T2D، گلوتامات و گیرنده‌های NMDA/AMPA ممکن است دچار دیس‌ریگولیشن شوند",
        "پلی‌فنول‌ها و امگا-3 می‌توانند مدولاسیون گیرنده‌ها/انتقال‌دهنده‌ها را تسهیل کنند (شواهد بیشتر پیش‌بالینی)"
    ])

    # 19
    add_bullets_slide("نتیجه تنظیم گلوتامات", [
        "کاهش بالقوه درد نوروپاتیک/تحریک‌سمّیت",
        "بهبود شناخت/خلق در مطالعات غیر T2D نیز گزارش شده (نیازمند شواهد انسانی مستقیم در T2D)"
    ])

    # 20
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "مکانیسم مکمل: تعدیل سیستم ایمنی")
    add_image(s, "assets/gut_immune.png", left=1.5, top=1.8, width=9.5)

    # 21
    add_bullets_slide("4) تعدیل میکروبیوم روده", [
        "افزایش تولید SCFA (بوتیرات) → GLP-1 ↑ و حساسیت به انسولین ↑",
        "افزایش غنا و گونه‌های مفید (Akkermansia, Faecalibacterium)"
    ])

    # 22
    add_bullets_slide("تأثیر بر متابولیسم سیستمیک", [
        "GLP-1 و PYY افزایش → کنترل اشتها/گلوکز",
        "سیگنالینگ اسیدهای صفراوی (FXR/TGR5) → گلوکز/چربی بهبود"
    ])

    # 23
    s = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(s, "مهار التهاب مزمن")
    add_image(s, "assets/inflammation.png", left=1.0, top=1.8, width=10.5)

    # 24
    add_bullets_slide("5) پتانسیل تسهیل رمیلیناسیون (برای عوارض عصبی T2D)", [
        "DHA/EPA و برخی پلی‌فنول‌ها می‌توانند میلین‌سازی را حمایت کنند (شواهد عمدتاً حیوانی/پیش‌بالینی)",
        "اهمیت در نوروپاتی دیابتی: حفاظت عصبی و ترمیم"
    ])

    # 25
    add_bullets_slide("اثر بر خلق و سلامت روانی", [
        "رژیم‌های ضدالتهاب (مانند مدیترانه‌ای) با بهبود خلق مرتبط‌اند → پایبندی بهتر، کنترل قند بهتر",
        "مدیریت استرس/خواب به‌عنوان عناصر سبک زندگی مکمل"
    ])

    # 26 — Sources table part 1 (مثل عکس شما)
    rows1 = [
      ["Esposito et al., 2009, Ann Intern Med",
       "RCT",
       "بیماران تازه‌تشخیص T2D؛ n≈215",
       "اثر رژیم مدیترانه‌ای بر کنترل قند و نیاز به دارو",
       "RCT دوگروهی؛ پیگیری تا 4 سال",
       "HbA1c کاهش بیشتر؛ تأخیر معنی‌دار در شروع دارو؛ کاهش وزن",
       "بهبود کنترل گلیسمی و کاهش نیاز دارو در T2D",
       "تک‌مرکزی؛ پایبندی خودگزارشی"],
      ["Shai et al., 2008, NEJM (DIRECT)",
       "RCT سه‌گروهی",
       "افراد دارای اضافه‌وزن؛ زیرگروه مبتلایان T2D",
       "مقایسه رژیم‌های کم‌چرب/مدیترانه‌ای/کم‌کربوهیدرات",
       "24 ماه؛ مداخله ساختاریافته محیط کاری",
       "در زیرگروه T2D، HbA1c و لیپیدها در مدیترانه‌ای بهتر از کم‌چرب",
       "مزیت مدیترانه‌ای در کنترل قند و پروفایل لیپیدی",
       "تحلیل زیرگروه؛ تعمیم‌پذیری محدود"],
    ]
    add_sources_table("جدول مقالات - بخش 1", rows1)

    # 27 — Sources table part 2
    rows2 = [
      ["Salas-Salvadó et al., 2011, Diabetes Care (PREDIMED-Reus)",
       "RCT پیشگیری",
       "افراد پرخطر غیر دیابتی؛ n≈418",
       "اثر رژیم مدیترانه‌ای بر بروز T2D",
       "گروه‌های با روغن زیتون/مغزدانه؛ پیگیری ≈4 سال",
       "کاهش ≈52% در بروز T2D نسبت به کنترل",
       "کاهش معنی‌دار ریسک بروز T2D",
       "زیرمطالعه؛ عوامل سبک زندگی"],
      ["Ajala et al., 2013, Am J Clin Nutr (Meta-analysis)",
       "مرور نظام‌مند و متاآنالیز",
       "بزرگسالان مبتلا به T2D در RCTها",
       "مقایسه رویکردهای رژیمی بر HbA1c/وزن/لیپید",
       "ترکیب نتایج چند مطالعه",
       "بیشترین کاهش HbA1c و بهبود وزن در مدیترانه‌ای",
       "برتری کلی مدیترانه‌ای در مدیریت T2D",
       "ناهمگنی مطالعات و تفاوت مداخلات"],
    ]
    add_sources_table("جدول مقالات - بخش 2", rows2)

    # 28
    add_bullets_slide("نتیجه‌گیری کلی از مطالعات", [
        "در مبتلایان: HbA1c ↓، حساسیت انسولین ↑، نیاز به دارو ↓، ریسک CVD ↓",
        "در پیشگیری: کاهش معنی‌دار خطر بروز T2D در جمعیت‌های پرخطر",
        "هم‌سویی مکانیسم‌های زیستی با نتایج بالینی",
        "قابلیت اجرا و پذیرش مناسب"
    ])

    # 29
    add_bullets_slide("پیشنهادات برای تحقیقات آینده", [
        "RCTهای سر به سر با GLP-1 RA/SGLT2i + مدیترانه‌ای",
        "بیومارکرهای مکانیسمی: میکروبیوم، متابولومیکس، سیگنالینگ صفراوی",
        "پیامدهای عصبی/شناختی در T2D",
        "راهبردهای ارتقای پایبندی و بومی‌سازی"
    ])

    # 30
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_shape(slide, "منابع و تشکر")
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.4))
    tf = tx.text_frame
    refs = [
        "Esposito K, et al. Ann Intern Med. 2009;151(5):306–314. doi:10.7326/0003-4819-151-5-200909010-00004",
        "Shai I, et al. N Engl J Med. 2008;359:229–241. doi:10.1056/NEJMoa0708681",
        "Salas-Salvadó J, et al. Diabetes Care. 2011;34(1):14–19. doi:10.2337/dc10-1288",
        "Ajala O, et al. Am J Clin Nutr. 2013;97(3):505–516. doi:10.3945/ajcn.112.042457",
        "نمودار شیوع: الهام از IDF Diabetes Atlas (کار آموزشی)"
    ]
    for i, r in enumerate(refs):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = r
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.RIGHT
    p = tf.add_paragraph()
    p.text = "با تشکر از توجه شما"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.RIGHT

    # Save
    prs.save("output/med-diet-t2d.pptx")

if __name__ == "__main__":
    build_deck()
